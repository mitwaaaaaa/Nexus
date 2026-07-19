import fitz  # PyMuPDF
import pdfplumber
import docx
import pptx
import pytesseract
from PIL import Image
import io
import os
import logging
from typing import List, Dict, Any, Tuple

logger = logging.getLogger(__name__)

class IngestionService:
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> Tuple[List[Dict[str, Any]], int]:
        chunks_by_page = []
        num_pages = 0
        
        try:
            # First try extracting metadata and counting pages using PyMuPDF
            doc = fitz.open(file_path)
            num_pages = len(doc)
            logger.info(f"PDF has {num_pages} pages. Starting text extraction.")
            
            # Using pdfplumber for high-quality text extraction
            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    page_num = page_idx + 1
                    text = page.extract_text()
                    
                    # If page has very little text, try OCR
                    if not text or len(text.strip()) < 100:
                        logger.info(f"Page {page_num} text density is low ({len(text) if text else 0} chars). Attempting OCR.")
                        try:
                            # Render page as image for OCR
                            fitz_page = doc.load_page(page_idx)
                            pix = fitz_page.get_pixmap(dpi=150)
                            img_data = pix.tobytes("png")
                            img = Image.open(io.BytesIO(img_data))
                            ocr_text = pytesseract.image_to_string(img)
                            if ocr_text and len(ocr_text.strip()) > len(text or ""):
                                text = ocr_text
                                logger.info(f"OCR successfully extracted {len(ocr_text)} chars from page {page_num}.")
                        except Exception as ocr_err:
                            logger.error(f"OCR failed for PDF page {page_num}: {ocr_err}")
                            
                    if text and text.strip():
                        chunks_by_page.append({
                            "page": page_num,
                            "text": text.strip()
                        })
            doc.close()
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            raise e
            
        return chunks_by_page, num_pages

    @staticmethod
    def extract_text_from_docx(file_path: str) -> Tuple[List[Dict[str, Any]], int]:
        chunks_by_page = []
        try:
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            
            combined_text = "\n".join(full_text)
            # Since word docs don't have standard "pages" in python-docx easily, we mock a page-based split
            # split by ~3000 chars per page
            page_size = 3000
            pages = [combined_text[i:i+page_size] for i in range(0, len(combined_text), page_size)]
            
            for idx, page_text in enumerate(pages):
                chunks_by_page.append({
                    "page": idx + 1,
                    "text": page_text.strip()
                })
            num_pages = len(pages)
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")
            raise e
        return chunks_by_page, num_pages

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> Tuple[List[Dict[str, Any]], int]:
        chunks_by_page = []
        try:
            prs = pptx.Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                combined_text = "\n".join(slide_text)
                if combined_text.strip():
                    chunks_by_page.append({
                        "page": idx + 1,
                        "text": combined_text.strip()
                    })
            num_pages = len(prs.slides)
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            raise e
        return chunks_by_page, num_pages

    @staticmethod
    def extract_text_from_txt(file_path: str) -> Tuple[List[Dict[str, Any]], int]:
        chunks_by_page = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            
            # Split into mock pages of ~3000 chars
            page_size = 3000
            pages = [text[i:i+page_size] for i in range(0, len(text), page_size)]
            for idx, page_text in enumerate(pages):
                chunks_by_page.append({
                    "page": idx + 1,
                    "text": page_text.strip()
                })
            num_pages = len(pages)
        except Exception as e:
            logger.error(f"Error reading TXT {file_path}: {e}")
            raise e
        return chunks_by_page, num_pages

    @staticmethod
    def extract_text_from_image(file_path: str) -> Tuple[List[Dict[str, Any]], int]:
        chunks_by_page = []
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            chunks_by_page.append({
                "page": 1,
                "text": text.strip()
            })
            num_pages = 1
        except Exception as e:
            logger.error(f"Error OCR-ing Image {file_path}: {e}")
            raise e
        return chunks_by_page, num_pages

    @classmethod
    def chunk_text(cls, extracted_pages: List[Dict[str, Any]], chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        chunks = []
        
        for page_data in extracted_pages:
            page_num = page_data["page"]
            text = page_data["text"]
            
            if len(text) <= chunk_size:
                chunks.append({
                    "page": page_num,
                    "text": text
                })
                continue
                
            start = 0
            while start < len(text):
                end = start + chunk_size
                chunk_txt = text[start:end]
                chunks.append({
                    "page": page_num,
                    "text": chunk_txt
                })
                start += (chunk_size - overlap)
                
        return chunks

    @classmethod
    def process_document(cls, file_path: str, file_type: str) -> Tuple[List[Dict[str, Any]], int]:
        file_type = file_type.lower()
        if file_type == "pdf":
            raw_pages, pages_count = cls.extract_text_from_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            raw_pages, pages_count = cls.extract_text_from_docx(file_path)
        elif file_type in ["pptx", "ppt"]:
            raw_pages, pages_count = cls.extract_text_from_pptx(file_path)
        elif file_type in ["png", "jpg", "jpeg", "webp", "tiff", "bmp"]:
            raw_pages, pages_count = cls.extract_text_from_image(file_path)
        else: # Default text file ingestion
            raw_pages, pages_count = cls.extract_text_from_txt(file_path)
            
        chunks = cls.chunk_text(raw_pages)
        return chunks, pages_count
